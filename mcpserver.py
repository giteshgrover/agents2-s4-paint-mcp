# basic import 
from mcp.server.fastmcp import FastMCP, Image
from mcp.server.fastmcp.prompts import base
from mcp.types import TextContent
from mcp.server.lowlevel import Server
from mcp import types
from PIL import Image as PILImage
import math
import sys
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os
import subprocess
from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
from google.auth.transport.requests import Request
from googleapiclient.discovery import build
from email.mime.text import MIMEText
import base64

# instantiate an MCP server client
mcp = FastMCP("Calculator")

# DEFINE TOOLS

#addition tool
@mcp.tool()
def add(a: int, b: int) -> int:
    """Add two numbers"""
    print("CALLED: add(a: int, b: int) -> int:")
    return int(a + b)

@mcp.tool()
def add_list(l: list) -> int:
    """Add all numbers in a list"""
    print("CALLED: add(l: list) -> int:")
    return sum(l)

# subtraction tool
@mcp.tool()
def subtract(a: int, b: int) -> int:
    """Subtract two numbers"""
    print("CALLED: subtract(a: int, b: int) -> int:")
    return int(a - b)

# multiplication tool
@mcp.tool()
def multiply(a: int, b: int) -> int:
    """Multiply two numbers"""
    print("CALLED: multiply(a: int, b: int) -> int:")
    return int(a * b)

#  division tool
@mcp.tool() 
def divide(a: int, b: int) -> float:
    """Divide two numbers"""
    print("CALLED: divide(a: int, b: int) -> float:")
    return float(a / b)

# power tool
@mcp.tool()
def power(a: int, b: int) -> int:
    """Power of two numbers"""
    print("CALLED: power(a: int, b: int) -> int:")
    return int(a ** b)

# square root tool
@mcp.tool()
def sqrt(a: int) -> float:
    """Square root of a number"""
    print("CALLED: sqrt(a: int) -> float:")
    return float(a ** 0.5)

# cube root tool
@mcp.tool()
def cbrt(a: int) -> float:
    """Cube root of a number"""
    print("CALLED: cbrt(a: int) -> float:")
    return float(a ** (1/3))

# factorial tool
@mcp.tool()
def factorial(a: int) -> int:
    """factorial of a number"""
    print("CALLED: factorial(a: int) -> int:")
    return int(math.factorial(a))

# log tool
@mcp.tool()
def log(a: int) -> float:
    """log of a number"""
    print("CALLED: log(a: int) -> float:")
    return float(math.log(a))

# remainder tool
@mcp.tool()
def remainder(a: int, b: int) -> int:
    """remainder of two numbers divison"""
    print("CALLED: remainder(a: int, b: int) -> int:")
    return int(a % b)

# sin tool
@mcp.tool()
def sin(a: int) -> float:
    """sin of a number"""
    print("CALLED: sin(a: int) -> float:")
    return float(math.sin(a))

# cos tool
@mcp.tool()
def cos(a: int) -> float:
    """cos of a number"""
    print("CALLED: cos(a: int) -> float:")
    return float(math.cos(a))

# tan tool
@mcp.tool()
def tan(a: int) -> float:
    """tan of a number"""
    print("CALLED: tan(a: int) -> float:")
    return float(math.tan(a))

# mine tool
@mcp.tool()
def mine(a: int, b: int) -> int:
    """special mining tool"""
    print("CALLED: mine(a: int, b: int) -> int:")
    return int(a - b - b)

@mcp.tool()
def create_thumbnail(image_path: str) -> Image:
    """Create a thumbnail from an image"""
    print("CALLED: create_thumbnail(image_path: str) -> Image:")
    img = PILImage.open(image_path)
    img.thumbnail((100, 100))
    return Image(data=img.tobytes(), format="png")

@mcp.tool()
def strings_to_chars_to_int(string: str) -> list[int]:
    """Return the ASCII values of the characters in a word"""
    print("CALLED: strings_to_chars_to_int(string: str) -> list[int]:")
    return [int(ord(char)) for char in string]

@mcp.tool()
def int_list_to_exponential_sum(int_list: list) -> float:
    """Return sum of exponentials of numbers in a list"""
    print("CALLED: int_list_to_exponential_sum(int_list: list) -> float:")
    return sum(math.exp(i) for i in int_list)

@mcp.tool()
def fibonacci_numbers(n: int) -> list:
    """Return the first n Fibonacci Numbers"""
    print("CALLED: fibonacci_numbers(n: int) -> list:")
    if n <= 0:
        return []
    fib_sequence = [0, 1]
    for _ in range(2, n):
        fib_sequence.append(fib_sequence[-1] + fib_sequence[-2])
    return fib_sequence[:n]

@mcp.tool()
def create_and_open_keynote_presentation() -> None:
    """Creates and opens a keynote presentation"""
    print("CALLED: create_and_open_keynote_presentation() -> None:")
    applescript = '''
    tell application "Keynote"
        activate
        if not (exists document 1) then
            set thisDoc to make new document with properties {document theme:theme "White"}
        end if
        tell the front document
            set thisSlide to slide 1
            -- Set slide layout to blank
            set base slide of thisSlide to master slide "Blank" of thisDoc
        end tell
    end tell
    '''
    try:
        result = subprocess.run(["osascript", "-e", applescript], 
                              capture_output=True, text=True, check=True)
        print(f"AppleScript executed successfully: {result.stdout}")
    except subprocess.CalledProcessError as e:
        print(f"AppleScript execution failed: {e.stderr}")
        raise
    except Exception as e:
        print(f"Unexpected error in AppleScript execution: {e}")
        raise

@mcp.tool()
def add_rectangle_in_keynote_presentation() -> None:
    """Adds a rectangle to an opened keynote presentation"""
    print("CALLED: add_rectangle_in_keynote_presentation() -> None:")
    applescript = '''
    tell application "Keynote"
        activate
        tell the front document
            set thisSlide to slide 1
            tell thisSlide
                set newShape to make new shape with properties {position:{150, 150}, width:300, height:300}
            end tell
        end tell
    end tell
    '''
    try:
        result = subprocess.run(["osascript", "-e", applescript], 
                              capture_output=True, text=True, check=True)
        print(f"AppleScript executed successfully: {result.stdout}")
    except subprocess.CalledProcessError as e:
        print(f"AppleScript execution failed: {e.stderr}")
        raise
    except Exception as e:
        print(f"Unexpected error in AppleScript execution: {e}")
        raise

@mcp.tool()
def add_text_in_keynote_presentation(text: str) -> None:
    """Adds text to an opened keynote presentation"""
    print("CALLED: add_text_in_keynote_presentation(text: str) -> None:")
    # Properly escape the text for AppleScript
    escaped_text = text.replace('"', '\\"').replace('\\', '\\\\')
    applescript = f'''
    tell application "Keynote"
        activate
        tell the front document
            set thisSlide to slide 1
            tell thisSlide
                set newShape to make new shape with properties {{object text:"{escaped_text}", position:{{150, 150}}, width:200, height:100}}
            end tell
        end tell
    end tell
    '''
    try:
        result = subprocess.run(["osascript", "-e", applescript], 
                              capture_output=True, text=True, check=True)
        print(f"AppleScript executed successfully: {result.stdout}")
    except subprocess.CalledProcessError as e:
        print(f"AppleScript execution failed: {e.stderr}")
        raise
    except Exception as e:
        print(f"Unexpected error in AppleScript execution: {e}")
        raise

# Define scopes for Gmail API access
# SCOPES = ['https://www.googleapis.com/auth/gmail.modify']
SCOPES = ["https://www.googleapis.com/auth/gmail.send"]

@mcp.tool()
def send_email(to_email: str, subject: str, body: str) -> None:
    """
    Send an email using Gmail API.

    Args:
        to_email (str): Recipient email address
        subject (str): Email subject
        body (str): Email body (plain text)
    """
        
    print("CALLED: send_email(to_email: str, subject: str, body: str) -> None:")
    creds = None

    # Load saved tokens if available
    credentials_path = "./.gmail-mcp/credentials.json"
    token_path = "./.gmail-mcp/token.json"
    if os.path.exists(token_path):
        creds = Credentials.from_authorized_user_file(token_path, SCOPES)
        print("token found")

     # If no valid creds, go through OAuth flow
    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            creds.refresh(Request())
        else:
            flow = InstalledAppFlow.from_client_secrets_file(credentials_path, SCOPES)
            creds = flow.run_local_server(port=0)
            print("Credential file found")
        # Save token for next time
        with open(token_path, "w") as token_file:
            token_file.write(creds.to_json())

    # Build Gmail service
    service = build("gmail", "v1", credentials=creds)

    # Create email message
    message = MIMEText(body)
    message["to"] = to_email
    message["from"] = "me"
    message["subject"] = subject
    raw_message = base64.urlsafe_b64encode(message.as_bytes()).decode()

    # Send the email
    send_result = service.users().messages().send(
        userId="me", body={"raw": raw_message}
    ).execute()
    print(f"Send Result: {send_result}")
    


def create_presentation() -> str:
    """Create a presentation and return the name of the presentation"""
    print("CALLED: create_presentation() -> str:")
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # blank slide (or prs.slide_layouts[1] ??)
    slide = prs.slides.add_slide(slide_layout)
    
    pr_name = "paint_like.pptx"
    prs.save(pr_name)
    print(f"Presentation saved as {pr_name}")
    # Open automatically (macOS)
    os.system(f"open {pr_name}")
    return pr_name


def add_rectangle_in_presentation(pr_name: str) -> None:
    """Add a rectangle to a presentation"""
    print("CALLED: add_rectangle(pr_name: str) -> None:")
    prs = Presentation(pr_name)
    slide = prs.slides[0]
    left = Inches(1)
    top = Inches(1)
    width = Inches(3)
    height = Inches(2)
    shape = slide.shapes.add_shape( 1,  # 1 = Rectangle (MSO_SHAPE.RECTANGLE), 
                                   left, top, width, height)
    # Format rectangle (fill + outline)
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 176, 80)  # green fill
    line = shape.line
    line.color.rgb = RGBColor(0, 0, 0)  # black border
    line.width = Pt(2)

    prs.save(pr_name)
    print(f"Presentation saved as {pr_name}")


def add_text_in_presentation(pr_name: str, text: str) -> None:
    """Add text to a presentation"""
    print("CALLED: add_text_in_presentation(pr_name: str, text: str) -> None:")
    prs = Presentation(pr_name)
    slide = prs.slides[0]
    txBox = slide.shapes.add_textbox(Inches(1.2), Inches(1.2), Inches(2.5), Inches(1))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)  # white text

    prs.save(pr_name)
    print(f"Presentation saved as {pr_name}")


if __name__ == "__main__":
    # Check if running with mcp dev command
    print("STARTING THE SERVER AT AMAZING LOCATION")
    if len(sys.argv) > 1 and sys.argv[1] == "dev":
        mcp.run()  # Run without transport for dev server
    else:
        mcp.run(transport="stdio")  # Run with stdio for direct execution
