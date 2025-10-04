import pyautogui
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os
import subprocess

# prs = Presentation()
# slide_layout = prs.slide_layouts[6]  # blank slide (or prs.slide_layouts[1] ??)
# slide = prs.slides.add_slide(slide_layout)

# pr_name = "paint_like.pptx"
# prs.save(pr_name)
# print(f"Presentation saved as {pr_name}")
# # Open automatically (macOS)
# os.system(f"open {pr_name}")

# time.sleep(3)  # give yourself 3s to switch to the drawing app

# # Draw a rectangle
# pyautogui.moveTo(200, 200)
# pyautogui.dragTo(400, 200, duration=0.5)
# pyautogui.dragTo(400, 400, duration=0.5)
# pyautogui.dragTo(200, 400, duration=0.5)
# pyautogui.dragTo(200, 200, duration=0.5)

# # Add text
# pyautogui.moveTo(250, 250)
# pyautogui.click()
# pyautogui.typewrite("Hello from Python!", interval=0.05)



applescript = '''
tell application "Keynote"
    activate
    if not (exists document 1) then
        set thisDoc to make new document with properties {document theme:theme "White"}
    end if
    tell the front document
        set thisSlide to slide 1
        -- Set slide layout to blank
        set base slide of thisSlide to master slide "Blank" of thisDoc
    end tell
end tell
'''

subprocess.run(["osascript", "-e", applescript])
time.sleep(3)

applescript = '''
tell application "Keynote"
    activate
    tell the front document
        set thisSlide to slide 1
        tell thisSlide
            set newShape to make new shape with properties {position:{150, 150}, width:300, height:300}
        end tell
    end tell
end tell
'''
subprocess.run(["osascript", "-e", applescript])
time.sleep(3)

applescript = '''
tell application "Keynote"
    activate
    tell the front document
        set thisSlide to slide 1
        tell thisSlide
            set newShape to make new shape with properties {object text:"Hello from Python!", position:{150, 150}, width:200, height:100}
        end tell
    end tell
end tell
'''
subprocess.run(["osascript", "-e", applescript])

